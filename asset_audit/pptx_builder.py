"""PPTX 生成器 - 三栏排版，参考抽卡分组大师样式"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# 幻灯片尺寸 16:9（标准宽屏）
SLIDE_W = Emu(9144000)
SLIDE_H = Emu(5143500)

# ---- 配色方案 ----
COLOR_TITLE    = RGBColor(0x0A, 0x0A, 0x0A)   # 主标题 近黑
COLOR_SUBTITLE = RGBColor(0x47, 0x55, 0x69)   # 副标题 深灰蓝
COLOR_ACCENT   = RGBColor(0x64, 0x74, 0x8B)   # 统计/标签 中灰蓝
COLOR_BODY     = RGBColor(0x33, 0x41, 0x55)   # 正文 深灰
COLOR_TAGLINE  = RGBColor(0x94, 0xA3, 0xB8)   # 页脚副文 浅灰蓝
COLOR_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_HEADING  = RGBColor(0x0A, 0x0A, 0x0A)   # 栏目标题 近黑
COLOR_BADGE_BG = RGBColor(0x7C, 0x4F, 0xD5)   # V2 徽章背景紫

# ---- 字体 ----
FONT_NAME = '微软雅黑'

# ---- 排版参数 ----
MARGIN        = Emu(457200)     # 左边距
CONTENT_TOP   = Emu(868680)     # 内容区 Y 起点
COL_WIDTH     = Emu(2697480)    # 单列宽
COL1_LEFT     = MARGIN          # 第 1 列 X
COL2_LEFT     = Emu(3246120)    # 第 2 列 X
COL3_LEFT     = Emu(6035040)    # 第 3 列 X

TITLE_SIZE       = Pt(24)       # 封面主标题
SUBTITLE_SIZE    = Pt(13)       # 封面副标题
STATS_SIZE       = Pt(10)       # 统计文字
SCENE_TITLE_SIZE = Pt(18)       # 场次/内容标题
COL_HEADING_SIZE = Pt(11)       # 栏目标题
BODY_SIZE        = Pt(8.5)      # 正文
FOOTER_TAG_SIZE  = Pt(7)        # 页脚标签
FOOTER_PAGE_SIZE = Pt(8)        # 页脚页码
BADGE_SIZE       = Pt(8)        # V2 徽章
APP_NAME_SIZE    = Pt(8.5)      # 应用名
HEADER_SIZE      = Pt(9)        # 页眉


def _set_font(run, size=BODY_SIZE, bold=None, color=COLOR_BODY):
    """统一设置字体属性"""
    run.font.name = FONT_NAME
    run.font.size = size
    if bold is not None:
        run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _add_textbox(slide, left, top, width, height):
    """创建空白文本框并返回 text_frame"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def _add_para(tf, text, size=BODY_SIZE, bold=None, color=COLOR_BODY,
              space_after=Pt(4), alignment=PP_ALIGN.LEFT):
    """向文本框添加段落"""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=bold, color=color)
    return p


def _add_footer(slide, page_num, total_pages):
    """添加统一页脚"""
    # V2 徽章
    badge = slide.shapes.add_shape(
        1, MARGIN, Emu(4709160), Emu(292608), Emu(201168))
    badge.fill.solid()
    badge.fill.fore_color.rgb = COLOR_BADGE_BG
    badge.line.fill.background()
    tf_badge = badge.text_frame
    tf_badge.word_wrap = False
    p = tf_badge.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = 'V2'
    _set_font(run, size=BADGE_SIZE, bold=True, color=COLOR_WHITE)

    # 应用名
    tf_app = _add_textbox(slide, Emu(822960), Emu(4645152), Emu(3200400), Emu(201168))
    _add_para(tf_app, '抽卡分组大师 Asset Master', size=APP_NAME_SIZE, bold=True, color=COLOR_TITLE)

    # 副标题
    tf_tag = _add_textbox(slide, Emu(822960), Emu(4828032), Emu(5486400), Emu(146304))
    _add_para(tf_tag, '智能剧本资产拆解与PPT排版系统 | Screenplay Asset Deconstruction & Grouping',
              size=FOOTER_TAG_SIZE, color=COLOR_TAGLINE)

    # 页码
    tf_page = _add_textbox(slide, COL3_LEFT, Emu(4736592), Emu(2651760), Emu(228600))
    _add_para(tf_page, f'第 {page_num} 页 / 共 {total_pages} 页',
              size=FOOTER_PAGE_SIZE, color=COLOR_TAGLINE, alignment=PP_ALIGN.RIGHT)


def build_episode_pptx(assets, series_name, episode_num, output_path):
    """为一集资产生成 PPTX - 三栏布局"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]
    ep_label = f'第{episode_num}集'

    chars = assets.get('characters', [])
    props_list = assets.get('props', [])
    scenes = assets.get('scenes', [])
    char_count = len(chars)
    prop_count = len(props_list)
    scene_count = len(scenes)

    # 计算内容页数
    MAX_PER_PAGE = 10
    content_pages = max(1, -(-max(char_count, prop_count, scene_count) // MAX_PER_PAGE))
    total_pages = 1 + content_pages

    # ========== 封面页 ==========
    slide = prs.slides.add_slide(blank_layout)

    tf = _add_textbox(slide, Emu(731520), Emu(1097280), Emu(7680960), Emu(1097280))
    _add_para(tf, f'{series_name} | {ep_label}', size=TITLE_SIZE, bold=True, color=COLOR_TITLE)

    tf = _add_textbox(slide, Emu(731520), Emu(2286000), Emu(7680960), Emu(548640))
    _add_para(tf, '单集剧本资产解构分组演示文稿', size=SUBTITLE_SIZE, color=COLOR_SUBTITLE)
    _add_para(tf, 'Screenplay Asset Deconstruction & Grouping Presentation',
              size=SUBTITLE_SIZE, color=COLOR_SUBTITLE)

    tf = _add_textbox(slide, Emu(731520), Emu(3017520), Emu(7680960), Emu(731520))
    _add_para(tf, '剧本资产统计 | Script Asset Summary:', size=STATS_SIZE, color=COLOR_ACCENT)
    stats = f'独立人物 | Characters: {char_count}  •  独立道具 | Props: {prop_count}  •  场景 | Scenes: {scene_count}'
    _add_para(tf, stats, size=STATS_SIZE, color=COLOR_ACCENT)

    _add_footer(slide, 1, total_pages)

    # ========== 内容页 ==========
    for page_idx in range(content_pages):
        slide = prs.slides.add_slide(blank_layout)
        page_num = page_idx + 2

        # 页眉
        tf = _add_textbox(slide, MARGIN, Emu(182880), Emu(8229600), Emu(228600))
        _add_para(tf, f'剧本 | Script: {series_name} | {ep_label}',
                  size=HEADER_SIZE, color=COLOR_ACCENT)

        # 内容标题
        start_n = page_idx * MAX_PER_PAGE + 1
        end_n = min((page_idx + 1) * MAX_PER_PAGE, max(char_count, prop_count, scene_count))
        tf = _add_textbox(slide, MARGIN, Emu(411480), Emu(8229600), Emu(365760))
        _add_para(tf, f'资产清单 | Asset List ({start_n}-{end_n})',
                  size=SCENE_TITLE_SIZE, bold=True, color=COLOR_TITLE)

        # 三栏标题
        tf_c1 = _add_textbox(slide, COL1_LEFT, CONTENT_TOP, COL_WIDTH, Emu(274320))
        _add_para(tf_c1, f'人物 | CHARACTERS（共{char_count}位）', size=COL_HEADING_SIZE, bold=True, color=COLOR_HEADING)

        tf_c2 = _add_textbox(slide, COL2_LEFT, CONTENT_TOP, COL_WIDTH, Emu(274320))
        _add_para(tf_c2, f'道具 | PROPS（共{prop_count}件）', size=COL_HEADING_SIZE, bold=True, color=COLOR_HEADING)

        tf_c3 = _add_textbox(slide, COL3_LEFT, CONTENT_TOP, Emu(2651760), Emu(274320))
        _add_para(tf_c3, f'场景 | SCENE（共{scene_count}个）', size=COL_HEADING_SIZE, bold=True, color=COLOR_HEADING)

        # 三栏内容
        col_y = Emu(1143000)

        # --- 人物 ---
        tf_chars = _add_textbox(slide, COL1_LEFT, col_y, COL_WIDTH, Emu(3657600))
        char_start_idx = page_idx * MAX_PER_PAGE
        for ci, c in enumerate(chars[page_idx * MAX_PER_PAGE:(page_idx + 1) * MAX_PER_PAGE]):
            name_cn = c.get('name_cn', '')
            name_en = c.get('name_en', '')
            costume = c.get('costume', '')
            line = f'{char_start_idx + ci + 1}. {name_cn}'
            if name_en:
                line += f' | {name_en}'
            _add_para(tf_chars, line, size=BODY_SIZE, color=COLOR_BODY)
            if costume:
                _add_para(tf_chars, f'  穿着：{costume}', size=BODY_SIZE,
                          color=COLOR_BODY, space_after=Pt(8))

        # --- 道具 ---
        tf_props = _add_textbox(slide, COL2_LEFT, col_y, COL_WIDTH, Emu(3657600))
        prop_start_idx = page_idx * MAX_PER_PAGE
        for pi, p in enumerate(props_list[page_idx * MAX_PER_PAGE:(page_idx + 1) * MAX_PER_PAGE]):
            name_cn = p.get('name_cn', '')
            usage = p.get('usage', '')
            _add_para(tf_props, f'{prop_start_idx + pi + 1}. {name_cn}', size=BODY_SIZE, color=COLOR_BODY)
            if usage:
                _add_para(tf_props, f'  情景：{usage}', size=BODY_SIZE,
                          color=COLOR_BODY, space_after=Pt(8))

        # --- 场景 ---
        tf_scenes = _add_textbox(slide, COL3_LEFT, col_y, Emu(2651760), Emu(3657600))
        scene_start_idx = page_idx * MAX_PER_PAGE
        for si, s in enumerate(scenes[page_idx * MAX_PER_PAGE:(page_idx + 1) * MAX_PER_PAGE]):
            name_cn = s.get('name_cn', '')
            name_en = s.get('name_en', '')
            synopsis = s.get('synopsis', '')
            scene_line = f'{scene_start_idx + si + 1}. 场景名称 : {name_cn}'
            if name_en:
                scene_line += f' | {name_en}'
            _add_para(tf_scenes, scene_line, size=BODY_SIZE, color=COLOR_BODY)
            if synopsis:
                _add_para(tf_scenes, f'场次概要 | Summary: {synopsis}', size=BODY_SIZE,
                          color=COLOR_BODY, space_after=Pt(8))

        _add_footer(slide, page_num, total_pages)

    os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else '.', exist_ok=True)
    prs.save(output_path)
    return output_path
